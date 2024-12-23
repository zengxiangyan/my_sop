# -*- coding: utf-8 -*-
import numpy as np
import xlsxwriter as xw
import openpyxl as op
import myoffice
import math
from openpyxl import Workbook,load_workbook
import re
from pptx import Presentation
from pptx.chart.xmlwriter import SeriesXmlRewriterFactory
from pptx.util import Inches, Pt
import pandas as pd
import time
import pythoncom
import os
import win32com
from myoffice.win32 import Win32
from win32com.client import constants, Dispatch, DispatchEx
from openpyxl.utils import get_column_letter,column_index_from_string
from pptx.dml.color import ColorFormat, RGBColor

def del_mydata(mydata):
    seasons = []
    delete = []
    for date in mydata['时间']:
        if date == '2024-01' or date == '2024-02' or date == '2024-03':
            seasons.append('2024 Q1')
        if date == '2024-04' or date == '2024-05' or date == '2024-06':
            seasons.append('2024 Q2')
        if date == '2024-07' or date == '2024-08' or date == '2024-09':
            seasons.append('2024 Q3')
        if date == '2024-10' or date == '2024-11' or date == '2024-12':
            seasons.append('2024 Q4')
        if date == '2023-01' or date == '2023-02' or date == '2023-03':
            seasons.append('2023 Q1')
        if date == '2023-04' or date == '2023-05' or date == '2023-06':
            seasons.append('2023 Q2')
        if date == '2023-07' or date == '2023-08' or date == '2023-09':
            seasons.append('2023 Q3')
        if date == '2023-10' or date == '2023-11' or date == '2023-12':
            seasons.append('2023 Q4')
        if date == '2022-01' or date == '2022-02' or date == '2022-03':
            seasons.append('2022 Q1')
        if date == '2022-04' or date == '2022-05' or date == '2022-06':
            seasons.append('2022 Q2')
        if date == '2022-07' or date == '2022-08' or date == '2022-09':
            seasons.append('2022 Q3')
        if date == '2022-10' or date == '2022-11' or date == '2022-12':
            seasons.append('2022 Q4')
        if date == '2021-01' or date == '2021-02' or date == '2021-03':
            seasons.append('2021 Q1')
        if date == '2021-04' or date == '2021-05' or date == '2021-06':
            seasons.append('2021 Q2')
        if date == '2021-07' or date == '2021-08' or date == '2021-09':
            seasons.append('2021 Q3')
        if date == '2021-10' or date == '2021-11' or date == '2021-12':
            seasons.append('2021 Q4')
        if date == '2020-01' or date == '2020-02' or date == '2020-03':
            seasons.append('2020 Q1')
        if date == '2020-04' or date == '2020-05' or date == '2020-06':
            seasons.append('2020 Q2')
        if date == '2020-07' or date == '2020-08' or date == '2020-09':
            seasons.append('2020 Q3')
        if date == '2020-10' or date == '2020-11' or date == '2020-12':
            seasons.append('2020 Q4')
    for name in mydata['名称']:
        if len(re.findall(r'星品|体验装|u先体验|试用装|体验礼', name)) > 0:
            delete.append('剔除')
        else:
            delete.append('否')
    mydata['seasons'] = seasons
    mydata['体验装剔除'] = delete
    return mydata


def sid_brand_datafill(mydata, table):
    w = 9
    sheet = table['total_panel']
    sid_sales = sid_sorted()  # 自然年累计
    sid_sales_fy = sid_sorted_fy_ytd()  # 财年累计
    sid_sales_recent_year = sid_sorted_recent_year()  # 过去一年
    # 用于计算22年度的店铺*季度数据
    season_total_sales = mydata.loc[mydata['seasons'].str.contains("|".join(date_range))].groupby(['seasons'], as_index=False)[['销售额', '去年同期销售额']].sum()
    season_total_sales = season_total_sales.copy()
    season_total_sales.reset_index(drop=True, inplace=True)
    lastyear_total_sales = mydata.loc[mydata['seasons'].str.contains("|".join(date_range))].groupby(['seasons'], as_index=False)['销售额'].sum()
    for i, season in enumerate(season_total_sales['seasons']):
        try:
            season_total_sales.loc[i, '去年同期销售额'] = lastyear_total_sales.loc[lastyear_total_sales['seasons'] == str(int(re.findall(r'\d*', season)[0]) - 1) + ' Q' + str(re.findall(r'\d*', season)[3])]['销售额'].values[0]
        except:
            season_total_sales.loc[i, '去年同期销售额'] = 0
    season_total_sales['YOY'] = (season_total_sales['销售额'] - season_total_sales['去年同期销售额']) / season_total_sales['去年同期销售额']
    for i in range(0, w):
        for j, season in enumerate(season_total_sales['seasons']):
            if season == sheet.cell(3 + 2 * i, 1).value:
                sheet.cell(3 + 2 * i, 2).value = round(season_total_sales['销售额'][j] / 1000000, 0)
                if math.isinf(float(season_total_sales['YOY'][j])):
                    sheet.cell(4 + 2 * i, 2).value = '-'  # 填充店铺面板第二列
                else:
                    sheet.cell(4 + 2 * i, 2).value = season_total_sales['YOY'][j]
    sheet.cell(3 + 2 * w, 2).value = round(season_total_sales.loc[season_total_sales['seasons'].str.contains("|".join(date_range3))]['销售额'].sum() / 1000000, 0)
    sheet.cell(4 + 2 * w, 2).value = (season_total_sales.loc[season_total_sales['seasons'].str.contains("|".join(date_range3))]['销售额'].sum() - season_total_sales.loc[season_total_sales['seasons'].str.contains("|".join(date_range3))]['去年同期销售额'].sum()) / season_total_sales.loc[season_total_sales['seasons'].str.contains("|".join(date_range3))]['去年同期销售额'].sum()  # 填充店铺面板第二列
    # 新增财年统计行
    sheet.cell(5 + 2 * w, 2).value = round(season_total_sales.loc[season_total_sales['seasons'].str.contains("|".join(date_range5))]['销售额'].sum() / 1000000, 0)
    sheet.cell(6 + 2 * w, 2).value = (season_total_sales.loc[season_total_sales['seasons'].str.contains("|".join(date_range5))]['销售额'].sum() - season_total_sales.loc[season_total_sales['seasons'].str.contains("|".join(date_range5))]['去年同期销售额'].sum()) / season_total_sales.loc[season_total_sales['seasons'].str.contains("|".join(date_range5))]['去年同期销售额'].sum()  # 填充店铺面板第二列
    # 新增过去一年统计行
    sheet.cell(7 + 2 * w, 2).value = round(season_total_sales.loc[season_total_sales['seasons'].str.contains("|".join(date_range7))]['销售额'].sum() / 1000000, 0)
    sheet.cell(8 + 2 * w, 2).value = (season_total_sales.loc[season_total_sales['seasons'].str.contains("|".join(date_range7))]['销售额'].sum() - season_total_sales.loc[season_total_sales['seasons'].str.contains("|".join(date_range7))]['去年同期销售额'].sum()) / season_total_sales.loc[season_total_sales['seasons'].str.contains("|".join(date_range7))]['去年同期销售额'].sum()  # 填充店铺面板第二列

    for num, sid in enumerate(sid_sales['店铺']):
        season_sales = mydata.loc[((mydata['店铺'] == sid) & (mydata['seasons'].str.contains("|".join(date_range))))].groupby(['seasons'], as_index=False)[['销售额', '去年同期销售额']].sum()
        lastyear_sales = mydata.loc[((mydata['店铺'] == sid) & (mydata['seasons'].str.contains("|".join(date_range))))].groupby(['seasons'], as_index=False)['销售额'].sum()
        season_sales = season_sales.copy()
        season_sales.reset_index(drop=True, inplace=True)
        for i, season in enumerate(season_sales['seasons']):
            try:
                season_sales.loc[i, '去年同期销售额'] = lastyear_sales.loc[lastyear_sales['seasons'] == str(int(re.findall(r'\d*', season)[0]) - 1) + ' Q' + str(re.findall(r'\d*', season)[3])]['销售额'].values[0]
            except:
                season_sales.loc[i, '去年同期销售额'] = 0
        season_sales['YOY'] = (season_sales['销售额'] - season_sales['去年同期销售额']) / season_sales['去年同期销售额']
        print("写入店铺*季度数据", sid)
        sheet.cell(1, num + 3).value = sid_names[sid]
        sheet.cell(2, num + 3).value = sid
        for row in range(w):
            for j, season in enumerate(season_sales['seasons']):
                if season == sheet.cell(3 + 2 * row, 1).value:
                    sheet.cell(3 + 2 * row, num + 3).value = round(season_sales['销售额'][j] / 1000000, 0)
                    if math.isinf(float(season_sales['YOY'][j])):
                        sheet.cell(3 + 2 * row + 1, num + 3).value = '-'
                    else:
                        sheet.cell(3 + 2 * row + 1, num + 3).value = season_sales['YOY'][j]
                    break
                else:
                    sheet.cell(3 + 2 * row, num + 3).value = 0
                    sheet.cell(3 + 2 * row + 1, num + 3).value = '-'
        sheet.cell(2 + 2 * w + 1, num + 3).value = round(list(sid_sales['销售额'])[num] / 1000000, 0)
        if math.isinf(float(list(sid_sales['YOY'])[num])):
            sheet.cell(2 + 2 * w + 1 + 1, num + 3).value = '-'
        else:
            sheet.cell(2 + 2 * w + 1 + 1, num + 3).value = list(sid_sales['YOY'])[num]
        # 财年分店铺统计行
        sheet.cell(2 + 2 * w + 1 + 2, num + 3).value = round(sid_sales_fy.loc[sid_sales_fy['店铺'] == sid]['销售额'].values[0] / 1000000, 0)
        if math.isinf(float(sid_sales_fy.loc[sid_sales_fy['店铺'] == sid]['YOY'].values[0])):
            sheet.cell(2 + 2 * w + 1 + 1 + 2, num + 3).value = '-'
        else:
            sheet.cell(2 + 2 * w + 1 + 1 + 2, num + 3).value = sid_sales_fy.loc[sid_sales_fy['店铺'] == sid]['YOY'].values[0]

        # 过去一年分店铺统计行
        sheet.cell(2 + 2 * w + 1 + 2 + 2, num + 3).value = round(
            sid_sales_recent_year.loc[sid_sales_recent_year['店铺'] == sid]['销售额'].values[0] / 1000000, 0)
        if math.isinf(float(sid_sales_recent_year.loc[sid_sales_recent_year['店铺'] == sid]['YOY'].values[0])):
            sheet.cell(2 + 2 * w + 1 + 1 + 2 + 2, num + 3).value = '-'
        else:
            sheet.cell(2 + 2 * w + 1 + 1 + 2 + 2, num + 3).value = sid_sales_recent_year.loc[sid_sales_recent_year['店铺'] == sid]['YOY'].values[0]
    return table


def function_brand_filldata(mydata, table):
    season_total_sales = seasons_panel()
    brand_sort, functions_sort = brand_sorted()
    total_sales = season_total_sales.loc[season_total_sales['seasons'].str.contains("|".join(date_range1))]['销售额'].sum()
    sheet = table['function_panel']
    brands_dict = {"kerastase/卡诗": "Kerastase*", "sisley": "Sisley", "欧舒丹": "L'occitane","馥绿德雅": "Rene Furterer", "L&rsquo;oreal professionnel/巴黎欧莱雅沙龙专属": "LorealPro*","grow gorgeous": "Grow Gorgeous", "SHISEIDO PROFESSIONAL/资生堂专业美发": "Shiseido Professional","MOROCCANOIL": "MOROCCANOIL*", "丰添": "Foltene*", "Living Proof": "Living Proof","Olaplex ": "Olaplex", "Christophe Robin": "Christophe Robin", "AVEDA": "AVEDA*","davines/大卫尼斯": "Davines"}

    # 填充第一列品牌名，按照品牌的销售额降序，不分功效
    for i, brand_name in enumerate(brand_sort):
        sheet.cell(3 + i, 1).value = brand_name
    # 填充第一行功效名称，按照功效的销售额降序，不分品牌
    for i, fun in enumerate(functions_sort):
        sheet.cell(1, 2 + i * 3).value = fun[0]
    for num, fun in enumerate(functions_sort):
        function = fun[0]
        brand_sales = mydata.loc[(mydata['seasons'].str.contains("|".join(date_range1))) & (mydata['功效-' + function] != '否')].groupby(['子品类', '品牌名', '功效-' + function], as_index=False)[['销售额', '去年同期销售额']].sum()
        brand_sales = brand_sales.copy()
        brand_sales.reset_index(drop=True, inplace=True)
        lastyear_brand_sales = mydata.loc[(mydata['seasons'].str.contains("|".join(date_range2))) & (mydata['功效-' + function] != '否')].groupby(['子品类', '品牌名', '功效-' + function], as_index=False)[['销售额', '去年同期销售额']].sum()
        function_sales = mydata.loc[(mydata['seasons'].str.contains("|".join(date_range1))) & (mydata['功效-' + function] != '否')].groupby(['子品类', '功效-' + function], as_index=False)[['销售额', '去年同期销售额']].sum()
        brand_share = []
        for i, brand in enumerate(brand_sales['品牌名']):
            try:
                brand_sales.loc[i, '去年同期销售额'] = lastyear_brand_sales.loc[lastyear_brand_sales['品牌名'] == brand]['销售额'].values[0]
            except:
                brand_sales.loc[i, '去年同期销售额'] = 0
            brand_share.append(list(brand_sales['销售额'])[i] / function_sales['销售额'][0])
        brand_sales['YOY'] = (brand_sales['销售额'] - brand_sales['去年同期销售额']) / brand_sales['去年同期销售额']
        brand_sales['brand_share'] = brand_share
        print("写入功效*店铺数据", function)
        for i in range(len(brand_sort)):
            for j, brand in enumerate(brand_sales['品牌名']):
                if brand == sheet.cell(3 + i, 1).value:
                    sheet.cell(3 + i, 3 * num + 2).value = round(brand_sales['销售额'][j] / 1000000, 2)
                    if math.isinf(float(brand_sales['YOY'][j])):
                        sheet.cell(3 + i, 3 * num + 3).value = '-'
                    else:
                        sheet.cell(3 + i, 3 * num + 3).value = brand_sales['YOY'][j]
                    sheet.cell(3 + i, 3 * num + 4).value = str(
                        round(round(brand_sales['brand_share'][j], 3) * 100, 1)) + '%'
                    break
                else:
                    sheet.cell(3 + i, 3 * num + 2).value = '-'
                    sheet.cell(3 + i, 3 * num + 3).value = '-'
                    sheet.cell(3 + i, 3 * num + 4).value = '-'
            if sheet.cell(3 + i, 3 * num + 3).value != '-':
                if sheet.cell(3 + i, 3 * num + 3).value > 0:
                    sheet.cell(3 + i, 3 * num + 3).value = '+' + str(int(round(sheet.cell(3 + i, 3 * num + 3).value, 2) * 100)) + '%'
        sheet.cell(3 + len(brand_sort), num * 3 + 2).value = round(fun[1] / 1000000, 0)
        sheet.cell(3 + len(brand_sort), num * 3 + 3).value = (fun[1] - brand_sales['去年同期销售额'].sum()) / brand_sales['去年同期销售额'].sum()
        sheet.cell(3 + len(brand_sort), num * 3 + 4).value = str(round(fun[1] / total_sales * 100, 1)) + '%'

    for i, brand_name in enumerate(brand_sort):
        sheet.cell(3 + i, 1).value = brands_dict[brand_name]
    return table

def function_top_item(mydata, table):
    functions = ['Nourish & repair', 'Anti-Hair Loss', 'Volumizing', 'Oil control', 'Color control', 'Anti-dandruff','Smooth & hydration', 'men care']
    brands_dict = {"kerastase/卡诗": "Kerastase*", "sisley": "Sisley", "欧舒丹": "L'occitane","馥绿德雅": "Rene Furterer", "L&rsquo;oreal professionnel/巴黎欧莱雅沙龙专属": "LorealPro*","grow gorgeous": "Grow Gorgeous", "SHISEIDO PROFESSIONAL/资生堂专业美发": "Shiseido Professional","MOROCCANOIL": "MOROCCANOIL*", "丰添": "Foltene*", "Living Proof": "Living Proof","Olaplex ": "Olaplex", "Christophe Robin": "Christophe Robin", "AVEDA": "AVEDA*","davines/大卫尼斯": "Davines"}
    for function in functions:
        item_sales = mydata.loc[(mydata['seasons'].str.contains("|".join(date_range1))) & (mydata['功效-' + function] != '否') & (mydata['体验装剔除'] != '剔除')].groupby(['子品类', '品牌名', '功效-' + function, 'tb_item_id'],as_index=False)[['销量', '销售额', '去年同期销售额']].sum().groupby(['子品类', '功效-' + function], as_index=False).apply(lambda x: x.nlargest(10, '销售额'))
        item_sales = item_sales.copy()
        item_sales.reset_index(drop=True, inplace=True)
        lastyear_item_sales = mydata.loc[(mydata['seasons'].str.contains("|".join(date_range2))) & (mydata['功效-' + function] != '否') & (mydata['体验装剔除'] != '剔除')].groupby(['子品类', '品牌名', '功效-' + function, 'tb_item_id'],as_index=False)[['销量', '销售额', '去年同期销售额']].sum()
        item_names = []
        trade_prop_values = []
        for i, item in enumerate(item_sales['tb_item_id']):
            item_name = mydata.loc[(mydata['seasons'].str.contains("|".join(date_range1))) & (mydata['功效-' + function] != '否') & (mydata['体验装剔除'] != '剔除')].groupby(['子品类', '品牌名', '功效-' + function, 'tb_item_id', '名称'], as_index=False)[['销量', '销售额', '去年同期销售额']].sum().groupby(['子品类', '功效-' + function], as_index=False).apply(lambda x: x.nlargest(200000, '销售额'))
            trade_prop_value = mydata.loc[(mydata['seasons'].str.contains("|".join(date_range1))) & (mydata['功效-' + function] != '否') & (mydata['体验装剔除'] != '剔除')].groupby(['子品类', '品牌名', '功效-' + function, 'tb_item_id', '名称', '交易属性'], as_index=False)[['销量', '销售额', '去年同期销售额']].sum().groupby(['子品类', '功效-' + function], as_index=False).apply(lambda x: x.nlargest(200000, '销售额'))
            if len(re.findall(r'[^【.*~】]+',list(item_name.loc[item_name['tb_item_id'] == item]['名称'])[0].replace('*', '').replace('.', '').replace('~', ''))) == 2:
                name = re.findall(r'[^【.*~】]+', list(item_name.loc[item_name['tb_item_id'] == item]['名称'])[0])[1]
            else:
                name = list(item_name.loc[item_name['tb_item_id'] == item]['名称'])[0]
            item_names.append(name)

            try:
                trade_prop = list(trade_prop_value.loc[trade_prop_value['tb_item_id'] == item]['交易属性'])[0]
            except:
                trade_prop = 'None'
            trade_prop_values.append(trade_prop)
            try:
                item_sales.loc[i, '去年同期销售额'] = lastyear_item_sales.loc[lastyear_item_sales['tb_item_id'] == item, '销售额'].values[0]
            except:
                item_sales.loc[i, '去年同期销售额'] = 0

        item_sales['YOY'] = (item_sales['销售额'] - item_sales['去年同期销售额']) / item_sales['去年同期销售额']
        item_sales['name'] = item_names
        item_sales['交易属性'] = trade_prop_values
        item_sales['Incremental'] = item_sales['销售额'] - item_sales['去年同期销售额']
        print("写入{}功效的top10宝贝".format(function))
        # 填充功效top10品牌名，按照品牌的销售额降序，分功效
        sheet = table[function]
        sheet.cell(2, 4).value = round(mydata.loc[(mydata['seasons'].str.contains("|".join(date_range2))) & (mydata['功效-' + function] != '否')]['销售额'].sum() / 1000000, 1)
        sheet.cell(2, 5).value = round(mydata.loc[(mydata['seasons'].str.contains("|".join(date_range1))) & (mydata['功效-' + function] != '否')]['销售额'].sum() / 1000000, 1)
        sheet.cell(2, 6).value = str(int(round((mydata.loc[(mydata['seasons'].str.contains("|".join(date_range1))) & (mydata['功效-' + function] != '否')]['销售额'].sum() - mydata.loc[(mydata['seasons'].str.contains("|".join(date_range2))) & (mydata['功效-' + function] != '否')]['销售额'].sum()) /mydata.loc[(mydata['seasons'].str.contains("|".join(date_range2))) & (mydata['功效-' + function] != '否')]['销售额'].sum() * 100,0))) + '%'
        sheet.cell(2, 7).value = int(round((mydata.loc[(mydata['seasons'].str.contains("|".join(date_range1))) & (mydata['功效-' + function] != '否')]['销售额'].sum() - mydata.loc[(mydata['seasons'].str.contains("|".join(date_range2))) & (mydata['功效-' + function] != '否')]['销售额'].sum()) / 1000000, 0))
        sheet.cell(2, 8).value = round(mydata.loc[(mydata['seasons'].str.contains("|".join(date_range1))) & (mydata['功效-' + function] != '否')]['销售额'].sum() /mydata.loc[(mydata['seasons'].str.contains("|".join(date_range1))) & (mydata['功效-' + function] != '否')]['销量'].sum(), 1)
        for i, brand_name in enumerate(list(item_sales['品牌名'])):
            sheet.cell(3 + i, 2).value = brands_dict[brand_name]
            sheet.cell(3 + i, 3).value = list(item_sales['name'])[i]
            if list(item_sales['销售额'])[i] / 1000000 >= 0.05:
                sheet.cell(3 + i, 4).value = round(list(item_sales['去年同期销售额'])[i] / 1000000, 1)
                sheet.cell(3 + i, 5).value = round(list(item_sales['销售额'])[i] / 1000000, 1)
            elif list(item_sales['销售额'])[i] / 1000000 >= 0.005:
                sheet.cell(3 + i, 4).value = round(list(item_sales['去年同期销售额'])[i] / 1000000, 2)
                sheet.cell(3 + i, 5).value = round(list(item_sales['销售额'])[i] / 1000000, 2)
            elif list(item_sales['销售额'])[i] / 1000000 >= 0.0005:
                sheet.cell(3 + i, 4).value = round(list(item_sales['去年同期销售额'])[i] / 1000000, 3)
                sheet.cell(3 + i, 5).value = round(list(item_sales['销售额'])[i] / 1000000, 3)
            else:
                sheet.cell(3 + i, 4).value = round(list(item_sales['去年同期销售额'])[i] / 1000000, 4)
                sheet.cell(3 + i, 5).value = round(list(item_sales['销售额'])[i] / 1000000, 4)
            sheet.cell(3 + i, 8).value = round(list(item_sales['销售额'])[i] / list(item_sales['销量'])[i], 1)
            if math.isinf(float(list(item_sales['YOY'])[i])):
                sheet.cell(3 + i, 4).value = '-'
                sheet.cell(3 + i, 6).value = '-'
            else:
                sheet.cell(3 + i, 6).value = list(item_sales['YOY'])[i]
            if sheet.cell(3 + i, 6).value == '-':
                sheet.cell(3 + i, 7).value = '+' + str(round(list(item_sales['销售额'])[i] / 1000000, 1))
            elif (list(item_sales['销售额'])[i] - list(item_sales['去年同期销售额'])[i]) >= 0:
                sheet.cell(3 + i, 7).value = '+' + str(
                    round((list(item_sales['销售额'])[i] - list(item_sales['去年同期销售额'])[i]) / 1000000, 1))
            else:
                sheet.cell(3 + i, 7).value = str(
                    round((list(item_sales['销售额'])[i] - list(item_sales['去年同期销售额'])[i]) / 1000000, 1))
            if function == 'men care':
                #                 sheet.cell(3+i,8).value = list(item_sales['交易属性'])[i]
                sheet.cell(3 + i, 3).value = list(item_sales['name'])[i] + '【' + list(item_sales['交易属性'])[i] + '】'

    return table

def table_updata_all(filename, filename2, slidenums=None, data= None, sheetnames=None):
    def fill_text_frame(frame, text):
        filled = False
        for paragraph in frame.paragraphs:
            if len(paragraph.runs) == 0:
                if not filled:
                    run = paragraph.add_run()
                    run.text = text
                    if paragraph.font.size:
                        run.font.size = paragraph.font.size
                    filled = True
            else:
                for run in paragraph.runs:
                    if not filled:
                        run.text = text
                        filled = True
                    else:
                        run.text = ''
        return filled
    if slidenums and len(sheetnames) == len(slidenums):
        # 新建ppt
        ppt = Presentation()
        # 新建页面
        for i in range(len(sheetnames)):
            slide = ppt.slides.add_slide(ppt.slide_layouts[6])
        # 保存ppt
        ppt.save(filename2)
        pythoncom.CoInitialize()  # excel多线程相关
        excel_object = DispatchEx("Excel.Application")
        excel_object.Visible = False
        excel_object.DisplayAlerts = False  # 是否显示警告
        excel_workbook = excel_object.Workbooks.Open(Filename=data)
        powerpoint_object = DispatchEx("Powerpoint.Application")
        powerpoint_object.Visible = True
        powerpoint_presentation = powerpoint_object.Presentations.Open(filename2)
        for i in range(len(sheetnames)):
            excel_worksheet = excel_workbook.Worksheets(sheetnames[i])
            excel_range = excel_worksheet.Range("A1:{}{}".format(get_column_letter(excel_worksheet.UsedRange.Columns.Count),excel_worksheet.UsedRange.Rows.Count))
            excel_range.Copy()
            time.sleep(1.5)
            powerpoint_presentation.Slides(i + 1).Shapes.PasteSpecial()# 使用 DataType=2 (ppPasteEnhancedMetafile) 保留格式
        powerpoint_presentation.Save()
        powerpoint_presentation.Close()
        powerpoint_object.Quit()
        excel_workbook.Close(False)
        excel_object.Quit()
        Win32.close(excel_workbook, powerpoint_presentation)

        myppt = Presentation(filename)
        myppt2 = Presentation(filename2)
        table = load_workbook(data)
        for i in range(len(slidenums)):
            slidenum = slidenums[i]
            sheet = table[sheetnames[i]]
            for shape in myppt.slides[slidenum].shapes:
                if shape.has_table:
                    if (len(shape.table.columns) == sheet.max_column) and (len(shape.table.rows) == sheet.max_row):
                        for col in range(len(shape.table.columns)):
                            for row in range(len(shape.table.rows)):
                                if hasattr(shape.table.cell(row, col), 'text_frame'):
                                    if not fill_text_frame(shape.table.cell(row, col).text_frame, myppt2.slides[i].shapes[0].table.cell(row, col).text):
                                        shape.table.cell.text = myppt2.slides[i].shapes[0].table.cell(row, col).text
                                else:
                                    shape.table.cell.text = myppt2.slides[i].shapes[0].table.cell(row, col).text
        myppt.save(filename)

def quarter_to_months(quarter_str):
    # 定义季度到英文月份的映射
    quarter_to_month_map = {
        'Q1': 'Jan-Mar',
        'Q2': 'Apr-Jun',
        'Q3': 'Jul-Sep',
        'Q4': 'Oct-Dec'
    }

    year, quarter = quarter_str.split()
    return f"{year} {quarter_to_month_map[quarter]}"


def convert_quarters_to_months(start_quarter, end_quarter):
    start_year, start_quarter = start_quarter.split()
    end_year, end_quarter = end_quarter.split()

    start_months = quarter_to_months(f"{start_year} {start_quarter}")
    end_months = quarter_to_months(f"{end_year} {end_quarter}")

    # 获取每个季度的起始和结束月份
    start_month = start_months.split()[1].split('-')[0]
    end_month = end_months.split()[1].split('-')[-1]

    return f"{start_month}-{end_month}"


def unit_table(table):
    w = 9
    sheet = table['total_panel']
    for i in range(0, w):
        sheet.cell(3 + 2 * i, 1).value = date_range[w-i-1]
    sheet.cell(3 + 2 * w, 1).value = convert_quarters_to_months(date_range3[len(date_range3)-1], date_range3[0])
    sheet.cell(3 + 2 * (w+2), 1).value = (date_range[3][2:] + '-' + date_range[0][2:]).replace(' ','')

    sheet = table['function_panel']
    sheet.cell(2, 1).value = date_range1[0]

    fun_sheet = ['Nourish & repair','Anti-Hair Loss','Volumizing','Oil control','Anti-dandruff','Smooth & hydration', 'Color control','men care']
    for fun in fun_sheet:
        sheet = table[fun]
        sheet.cell(1, 4).value = date_range2[0]
        sheet.cell(1, 5).value = date_range1[0]

    return table

##tortal_panel最后两行数据
def sid_sorted():
    sid_sales = mydata.loc[mydata['seasons'].str.contains("|".join(date_range3))].groupby(['子品类', '店铺'], as_index=False)[['销售额', '去年同期销售额']].sum()
    sid_sales = sid_sales.copy()
    sid_sales.reset_index(drop=True, inplace=True)
    for sid in sid_names.keys():
        if sid not in list(sid_sales['店铺']):
            sid_sales = sid_sales.append({'子品类': '美发护发', '店铺': sid, '销售额': 0, '去年同期销售额': 0},ignore_index=True)
    lastyear_sid_sales = mydata.loc[mydata['seasons'].str.contains("|".join(date_range4))].groupby(['子品类', '店铺'], as_index=False)[['销售额', '去年同期销售额']].sum()
    for i, sid in enumerate(sid_sales['店铺']):
        try:
            sid_sales.loc[i, '去年同期销售额'] = lastyear_sid_sales.loc[lastyear_sid_sales['店铺'] == sid]['销售额'].values[0]
        except:
            sid_sales.loc[i, '去年同期销售额'] = 0.0
    sid_sales['YOY'] = (sid_sales['销售额'] - sid_sales['去年同期销售额']) / sid_sales['去年同期销售额']
    sid_sales = sid_sales.groupby(['子品类'], as_index=False).apply(lambda x: x.nlargest(20000, '销售额'))
    return sid_sales


def sid_sorted_fy_ytd():
    sid_sales_fy = mydata.loc[mydata['seasons'].str.contains("|".join(date_range5))].groupby(['子品类', '店铺'], as_index=False)[['销售额', '去年同期销售额']].sum()
    sid_sales_fy = sid_sales_fy.copy()
    sid_sales_fy.reset_index(drop=True, inplace=True)
    for sid in sid_names.keys():
        if sid not in list(sid_sales_fy['店铺']):
            sid_sales_fy = sid_sales_fy.append({'子品类': '美发护发', '店铺': sid, '销售额': 0, '去年同期销售额': 0},
                                               ignore_index=True)
    lastyear_sid_sales_fy = mydata.loc[mydata['seasons'].str.contains("|".join(date_range6))].groupby(['子品类', '店铺'], as_index=False)[['销售额', '去年同期销售额']].sum()
    for i, sid in enumerate(sid_sales_fy['店铺']):
        try:
            sid_sales_fy.loc[i, '去年同期销售额'] = lastyear_sid_sales_fy.loc[lastyear_sid_sales_fy['店铺'] == sid]['销售额'].values[0]
        except:
            sid_sales_fy.loc[i, '去年同期销售额'] = 0.0
    sid_sales_fy['YOY'] = (sid_sales_fy['销售额'] - sid_sales_fy['去年同期销售额']) / sid_sales_fy['去年同期销售额']
    sid_sales_fy = sid_sales_fy.groupby(['子品类'], as_index=False).apply(lambda x: x.nlargest(20000, '销售额'))
    return sid_sales_fy


def sid_sorted_recent_year():
    sid_sales_recent_year = mydata.loc[mydata['seasons'].str.contains("|".join(date_range7))].groupby(['子品类', '店铺'], as_index=False)[['销售额', '去年同期销售额']].sum()
    sid_sales_recent_year = sid_sales_recent_year.copy()
    sid_sales_recent_year.reset_index(drop=True, inplace=True)
    for sid in sid_names.keys():
        if sid not in list(sid_sales_recent_year['店铺']):
            sid_sales_recent_year = sid_sales_recent_year.append({'子品类': '美发护发', '店铺': sid, '销售额': 0, '去年同期销售额': 0}, ignore_index=True)
    lastyear_sid_sales_recent_year = mydata.loc[mydata['seasons'].str.contains("|".join(date_range8))].groupby(['子品类', '店铺'], as_index=False)[['销售额', '去年同期销售额']].sum()
    for i, sid in enumerate(sid_sales_recent_year['店铺']):
        try:
            sid_sales_recent_year.loc[i, '去年同期销售额'] = lastyear_sid_sales_recent_year.loc[lastyear_sid_sales_recent_year['店铺'] == sid]['销售额'].values[0]
        except:
            sid_sales_recent_year.loc[i, '去年同期销售额'] = 0.0
    sid_sales_recent_year['YOY'] = (sid_sales_recent_year['销售额'] - sid_sales_recent_year['去年同期销售额']) / sid_sales_recent_year['去年同期销售额']
    sid_sales_recent_year = sid_sales_recent_year.groupby(['子品类'], as_index=False).apply(lambda x: x.nlargest(20000, '销售额'))
    return sid_sales_recent_year


def seasons_panel():
    season_total_sales = mydata.loc[mydata['seasons'].str.contains("|".join(date_range))].groupby(['seasons'], as_index=False)[['销售额', '去年同期销售额']].sum()
    lastyear_total_sales = mydata.loc[mydata['seasons'].str.contains("|".join(date_range))].groupby(['seasons'], as_index=False)['销售额'].sum()
    season_total_sales = season_total_sales.copy()
    season_total_sales.reset_index(drop=True, inplace=True)
    for i, season in enumerate(season_total_sales['seasons']):
        try:
            season_total_sales.loc[i, '去年同期销售额'] = lastyear_total_sales.loc[lastyear_total_sales['seasons'] == str(int(re.findall(r'\d*', season)[0]) - 1) + ' Q' + str(re.findall(r'\d*', season)[3])]['销售额'].values[0]
        except:
            season_total_sales.loc[i, '去年同期销售额'] = 0
    season_total_sales['YOY'] = (season_total_sales['销售额'] - season_total_sales['去年同期销售额']) / season_total_sales['去年同期销售额']
    return season_total_sales


def brand_sorted():
    functions = ['Nourish & repair', 'Anti-Hair Loss', 'Volumizing', 'Oil control', 'Color control', 'Anti-dandruff','Smooth & hydration', 'men care']
    functions_sort = {}
    for function in functions:
        function_sales = mydata.loc[(mydata['seasons'].str.contains("|".join(date_range1))) & (mydata['功效-' + function] != '否')].groupby(['子品类', '功效-' + function], as_index=False)[['销售额', '去年同期销售额']].sum()
        lastyear_function_sales = mydata.loc[(mydata['seasons'].str.contains("|".join(date_range2))) & (mydata['功效-' + function] != '否')].groupby(['子品类', '功效-' + function], as_index=False)[['销售额', '去年同期销售额']].sum()
        functions_sort[function] = function_sales['销售额'][0]
    functions_sort = sorted(functions_sort.items(), key=lambda x: x[1], reverse=True)
    brand_sales = mydata.loc[mydata['seasons'].str.contains("|".join(date_range1))].groupby(['子品类', '品牌名'], as_index=False)[['销售额', '去年同期销售额']].sum().groupby(['子品类'], as_index=False).apply(lambda x: x.nlargest(2000, '销售额'))
    brand_sort = list(brand_sales['品牌名'])
    return brand_sort, functions_sort


def open_table(filename):
    """
    filename:C:\\Users\\asus\\Downloads\\PPT数据表.xlsx
    table:Using openpyxl open filename as table

    """
    filename = filename
    table = load_workbook(filename)

    table = unit_table(table)
    return table


def save_table(filename, table):
    '''
    filename:C:\\Users\\asus\\Downloads\\PPT数据表(2.8).xlsx
    table:the table of having been writed

    '''
    filename = filename
    table.save(filename)
    return table

def update_pptx(savepath,output_ppt,update_data=True):
    template = r'C:\Users\zeng.xiangyan\雅诗兰黛头发功效(指定店铺)\Prestige haircare panel_2023Q1 v2.pptx'
    filename = r'C:\Users\zeng.xiangyan\雅诗兰黛头发功效(指定店铺)\Prestige haircare panel_2022 Jan-Jul(2.0).pptx'
    filename2 = r'C:\Users\zeng.xiangyan\雅诗兰黛头发功效(指定店铺)\OUTPUT(2.0).pptx'

    # 以 template为模板，创建 output.pptx 文件,  填入datafile.xlsx 文件数据, 保存
    if update_data:
        myoffice.open_file(filename2,template=template).fill(savepath).save()
        table_updata_all(filename2, filename, slidenums=[0,2,3,4,5,6,7,8,9,10], data= savepath,
            sheetnames=['total_panel','function_panel',
                        'Nourish & repair',
                         'Anti-Hair Loss',
                         'Volumizing',
                         'Oil control',
                         'Anti-dandruff',
                         'Smooth & hydration',
                         'Color control',
                         'men care'])

    # 设置pptx功效品牌面板的YOY字体颜色
    myppt = Presentation(filename2)
    for shape in myppt.slides[2].shapes:
        if shape.has_table:
            for row in range(len(shape.table.rows) - 1):
                for col in range(len(shape.table.columns)):
                    if shape.table.cell(row, col).text_frame.text != '-':
                        try:
                            if len(re.findall(r'[+%]', shape.table.cell(row, col).text_frame.text)) == 2:
                                shape.table.cell(row, col).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0,102,0)
                            if len(re.findall(r'[-%]', shape.table.cell(row, col).text_frame.text)) == 2:
                                shape.table.cell(row, col).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(192, 0, 0)
                        except:
                            continue

    # 设置pptx total_panel页三个汇总行的条件格式
    index_list = ['ytd','fy ytd','mat']
    for shape in myppt.slides[0].shapes:
        if shape.has_table:
            sheet_table1 = shape.table
    for shape in myppt.slides[1].shapes:
        if shape.has_table:
            sheet_table2 = shape.table
            print(len(sheet_table1.rows), len(sheet_table2.rows))
    for id,index in enumerate(index_list):
        for i in range(0, len(sheet_table2.columns)):
            try:
                if sheet_table2.cell(len(sheet_table2.rows) - 1-(id*2), i).fill.fore_color.rgb:
                    sheet_table1.cell(len(sheet_table2.rows) - 1-(id*2), i).fill.patterned()
                    sheet_table1.cell(len(sheet_table2.rows) - 1-(id*2), i).fill.fore_color.rgb = sheet_table2.cell(len(sheet_table2.rows) - 1-(id*2), i).fill.fore_color.rgb
                    sheet_table1.cell(len(sheet_table2.rows) - 1-(id*2), i).fill.back_color.rgb = sheet_table2.cell(len(sheet_table2.rows) - 1-(id*2), i).fill.fore_color.rgb
            except:
                continue

    myppt.save(output_ppt)
def run(data, filepath, savepath,output_ppt):
    '''
    data:      the data is from sop_e.entify_prod_92111_E
    filepath:  the data is a model,being used for filling datas
    savepath:  the data of being filled datas

    '''

    global mydata, table, sid_names, date_range, date_range1, date_range2, date_range3, date_range4, date_range5, date_range6, date_range7, date_range8
    sid_names = {"kerastase卡诗官方旗舰店": "KERASTASE", "SISLEY希思黎官方旗舰店": "Sisley",
                 "L’OCCITANE欧舒丹官方旗舰店": "L'occitane", "馥绿德雅旗舰店": "RENE FURTERER",
                 "LorealPro官方旗舰店": "LorealPro", "资生堂专业美发官方旗舰店": "SHISEIDO PROFESSIONAL",
                 "OLAPLEX海外旗舰店": "OLAPLEX", "GrowGorgeous海外旗舰店": "GROW GORGEOUS",
                 "MOROCCANOIL海外旗舰店": "Moroccanoil", "丰添旗舰店": "Foltene",
                 "davines大卫尼斯海外旗舰店": "Davines", "LivingProof海外旗舰店": "Living proof",
                 "ChristopheRobin海外旗舰店": "Christophe Robin", "kerastase卡诗海外旗舰店": "KERASTASE Overseas Store",
                 "摩洛哥油旗舰店": "Moroccanoil", "LorealPro海外旗舰店": "LorealPro Overseas Store",
                 "AVEDA艾梵达官方旗舰店": "Aveda", "Aveda官方海外旗舰店": "Aveda Overseas Store"}
    mydata = pd.read_csv(data, encoding='utf-8')
    date_range = ['2024 Q3','2024 Q2', '2024 Q1', '2023 Q4', '2023 Q3', '2023 Q2', '2023 Q1', '2022 Q4', '2022 Q3', '2022 Q2',
                  '2022 Q1', '2021 Q4', '2021 Q3', '2021 Q2', '2021 Q1', '2020 Q4', '2020 Q3', '2020 Q2', '2020 Q1']
    # date_range1 = ['2024 Q3']  # 本期报告的数据时间范围(ppt第二页时间范围，常规为对应季度)
    # date_range2 = ['2023 Q3']  # 本期报告的数据时间范围(ppt第二页时间范围，常规为对应季度)，取历史数据计算同比
    date_range1 = ['2024 Q1','2024 Q2','2024 Q3']#临时改时间段出报告
    date_range2 = ['2023 Q1','2023 Q2','2023 Q3']
    date_range3 = ['2024 Q3','2024 Q2', '2024 Q1']  # 自然年统计时间范围
    date_range4 = ['2023 Q3','2023 Q2', '2023 Q1']
    date_range5 = ['2024 Q3']  # 财年从7月开始
    date_range6 = ['2023 Q3']
    date_range7 = ['2024 Q3','2024 Q2', '2024 Q1', '2023 Q4']  # 过去一年的统计
    date_range8 = ['2023 Q3','2023 Q2', '2023 Q1', '2022 Q4']

    del_mydata(mydata)
    table = open_table(filepath)
    sid_brand_datafill(mydata, table)
    function_brand_filldata(mydata, table)
    function_top_item(mydata, table)
    save_table(savepath, table)
    print("数据更新完毕，请去'{}'下检查！".format(savepath))

    update_pptx(savepath,output_ppt)

def main():
    # run('C:\\Users\\zeng.xiangyan\\Downloads\\92111 (1).csv',
    #     'C:\\Users\\zeng.xiangyan\\雅诗兰黛头发功效(指定店铺)\\PPT数据表(2.8).xlsx',
    #     'C:\\Users\\zeng.xiangyan\\雅诗兰黛头发功效(指定店铺)\\PPT数据表_24Q1-Q3（check）v1.xlsx',
    #     'C:\\Users\\zeng.xiangyan\\雅诗兰黛头发功效(指定店铺)\\Prestige haircare panel_24Q1-Q3(12-10).pptx')

    update_pptx('C:\\Users\\zeng.xiangyan\\雅诗兰黛头发功效(指定店铺)\\PPT数据表_24Q1-Q3（check）v1.xlsx',
        'C:\\Users\\zeng.xiangyan\\雅诗兰黛头发功效(指定店铺)\\Prestige haircare panel_24Q1-Q3(12-10).pptx',update_data=False)
if __name__ == "__main__":
    main()
    # unit_table(table=open_table('C:\\Users\\zeng.xiangyan\\雅诗兰黛头发功效(指定店铺)\\PPT数据表(2.8).xlsx'))
# sid_sales = mydata.loc[(mydata['seasons']=='2021Q1') | (mydata['seasons']=='2021Q2') | (mydata['seasons']=='2021Q3')].groupby(['子品类','店铺'],as_index=False)['销售额','去年同期销售额'].sum()
# sid_names = {"卡诗官方旗舰店":"KERASTASE","sisley希思黎官方旗舰店":"Sisley","loccitane欧舒丹旗舰店":"L'occitane","馥绿德雅旗舰店":"RENE FURTERER","lorealpro官方旗舰店":"LorealPro","资生堂专业美发官方旗舰店":"SHISEIDO PROFESSIONAL","OLAPLEX海外旗舰店":"OLAPLEX","GrowGorgeous海外旗舰店":"GROW GORGEOUS","moroccanoil海外旗舰店":"Moroccanoil","丰添旗舰店":"Foltene","davines大卫尼斯海外旗舰店":"Davines","LivingProof海外旗舰店":"Living proof","ChristopheRobin海外旗舰":"Christophe Robin","kerastase卡诗海外旗舰店":"KERASTASE Overseas Store","摩洛哥油旗舰店":"Moroccanoil","LorealPro海外旗舰店":"LorealPro Overseas Store"}
# sid_sales['YOY'] = (sid_sales['销售额']-sid_sales['去年同期销售额'])/sid_sales['去年同期销售额']season_sales
# sid_sales